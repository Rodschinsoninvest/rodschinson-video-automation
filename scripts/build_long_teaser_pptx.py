"""Build an EDITABLE PowerPoint version of a long teaser.

Reads the same JSON consumed by puppeteer/templates/teaser_long.html and
emits a .pptx where every section is rebuilt as native PowerPoint shapes
(rectangles + text boxes + pictures) so the user can edit text, swap
images, reflow tables, etc. directly in PowerPoint.

Usage:
  python scripts/build_long_teaser_pptx.py \
      --script output/teaser/e363e0b0_long_teaser.json \
      --output  output/teaser/e363e0b0_long_teaser.pptx
"""
from __future__ import annotations

import argparse
import json
import os
from urllib.parse import unquote, urlparse

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ─── Brand palette ─────────────────────────────────────────────────────────
NAVY      = RGBColor(0x08, 0x31, 0x6F)
NAVY_DEEP = RGBColor(0x06, 0x1F, 0x47)
SKY       = RGBColor(0x00, 0xB6, 0xFF)
GOLD      = RGBColor(0xC8, 0xA9, 0x6E)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREY_BG   = RGBColor(0xF5, 0xF7, 0xFA)
GREY_LINE = RGBColor(0xE5, 0xE7, 0xEB)
GREY_TX   = RGBColor(0x6B, 0x72, 0x80)
DARK_TX   = RGBColor(0x11, 0x18, 0x27)

# ─── Slide geometry (16:9 widescreen, inches) ──────────────────────────────
SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN_X = 0.45
HEADER_H = 0.85
BODY_TOP = 1.05
BODY_BOTTOM = SLIDE_H - 0.35


# ─── Helpers ───────────────────────────────────────────────────────────────
def file_path_from_url(u: str) -> str:
    """Accept file:// URLs or plain paths and return a local FS path."""
    if not u:
        return ''
    if u.startswith('file://'):
        return unquote(urlparse(u).path)
    return u


def add_rect(slide, x, y, w, h, fill=None, line=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.shadow.inherit = False
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        if line_width is not None:
            shape.line.width = Pt(line_width)
    return shape


def add_text(slide, x, y, w, h, text, *, size=11, bold=False, color=DARK_TX,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font='Calibri',
             letter_spacing=None, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top  = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = str(text).split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        if letter_spacing is not None:
            # XML kerning hack: spc attribute in 100ths of a point
            from pptx.oxml.ns import qn
            rPr = run._r.get_or_add_rPr()
            rPr.set('spc', str(int(letter_spacing * 100)))
    return tb


def add_picture(slide, path, x, y, w, h, *, crop_object_fit_cover=True):
    """Insert a picture with optional CSS-cover-like cropping so it fills
    the given frame without distortion."""
    if not path or not os.path.exists(path):
        # leave a placeholder rectangle in light grey
        add_rect(slide, x, y, w, h, fill=GREY_BG, line=GREY_LINE)
        return None
    pic = slide.shapes.add_picture(path, Inches(x), Inches(y),
                                   width=Inches(w), height=Inches(h))
    if not crop_object_fit_cover:
        return pic
    # Compute aspect ratios and apply symmetric crops (in % of original dimension)
    try:
        from PIL import Image
        with Image.open(path) as im:
            src_w, src_h = im.size
        src_ratio = src_w / src_h
        box_ratio = w / h
        if abs(src_ratio - box_ratio) < 0.01:
            return pic
        if src_ratio > box_ratio:
            # source wider — crop left/right
            new_w = src_h * box_ratio
            crop_x = (src_w - new_w) / 2 / src_w  # fraction
            pic.crop_left = pic.crop_right = crop_x
        else:
            # source taller — crop top/bottom
            new_h = src_w / box_ratio
            crop_y = (src_h - new_h) / 2 / src_h
            pic.crop_top = pic.crop_bottom = crop_y
    except Exception:
        pass
    return pic


def page_chrome(slide, data, nav_active, section_label):
    """Draw the top blue-gradient header strip + nav tabs row that
    appears on every interior slide (Activa, Details, Localisation, etc.)."""
    # Header bar
    bar = add_rect(slide, 0, 0, SLIDE_W, HEADER_H,
                   fill=NAVY)
    # Sky gradient overlay (simulated as two stacked rectangles since pptx
    # doesn't allow easy linear gradients via the API). Use one for now.
    # Wordmark left
    add_text(slide, MARGIN_X, 0.18, 3.8, 0.32, 'R O D S C H I N S O N',
             size=11.5, bold=False, color=WHITE, letter_spacing=4)
    add_text(slide, MARGIN_X, 0.52, 3.8, 0.22, 'I N V E S T M E N T',
             size=7.5, color=SKY, letter_spacing=5)
    # Section label centre-ish
    add_text(slide, 4.6, 0.13, 5.0, 0.34, section_label,
             size=14, bold=True, color=WHITE)
    loc = (data.get('title') or '').split(':', 1)[0].strip()
    add_text(slide, 4.6, 0.46, 5.0, 0.28, loc,
             size=8, color=WHITE, letter_spacing=2)
    add_text(slide, 4.6, 0.46, 5.0, 0.28, ' · ' + (data.get('reference','#REF') or ''),
             size=8, color=WHITE)
    # ref column on the right
    add_text(slide, 9.5, 0.46, 3.0, 0.28, data.get('reference',''),
             size=8.5, color=WHITE, letter_spacing=2)

    # Nav tab row right under the header
    tabs_y = HEADER_H + 0.05
    add_rect(slide, 0, tabs_y, SLIDE_W, 0.5, fill=WHITE,
             line=GREY_LINE, line_width=0.5)
    labels = [
        ('a',  data.get('tab_activa', 'Actif')),
        ('l',  data.get('tab_locatie', 'Localisation')),
        ('ae', data.get('tab_aerial', 'Vue aérienne')),
        ('p',  data.get('tab_photos', 'Photos')),
        ('s',  data.get('tab_sales',  'Conditions de vente')),
    ]
    x = MARGIN_X + 0.4
    for key, label in labels:
        is_active = key == nav_active
        # text width estimate
        w = max(0.85, 0.11 * len(label) + 0.4)
        add_text(slide, x, tabs_y + 0.12, w, 0.28, label,
                 size=10, bold=is_active,
                 color=NAVY if is_active else GREY_TX)
        # separator pipe
        x += w
        add_text(slide, x, tabs_y + 0.12, 0.15, 0.28, '|',
                 size=10, color=GREY_LINE)
        x += 0.15


def add_section_title(slide, x, y, w, h, title, *, suffix=''):
    add_text(slide, x, y, w - 1.5, h, title,
             size=18, bold=True, color=NAVY)
    if suffix:
        add_text(slide, x + w - 1.5, y + 0.05, 1.5, h, suffix,
                 size=9, color=GREY_TX, align=PP_ALIGN.RIGHT,
                 letter_spacing=2)
    # underline
    add_rect(slide, x, y + h, w, 0.012, fill=GREY_LINE)


# ─── Per-slide builders ────────────────────────────────────────────────────
def build_cover(slide, data):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY_DEEP)
    # reference top right
    add_text(slide, SLIDE_W - 2.4, 0.45, 2.0, 0.3, data.get('reference', ''),
             size=10, color=SKY, align=PP_ALIGN.RIGHT, letter_spacing=3)
    # Logo
    add_text(slide, SLIDE_W / 2 - 3, 2.3, 6, 0.6, 'R O D S C H I N S O N',
             size=28, color=WHITE, align=PP_ALIGN.CENTER, letter_spacing=8)
    add_rect(slide, SLIDE_W / 2 - 0.9, 2.92, 1.8, 0.025, fill=SKY)
    add_text(slide, SLIDE_W / 2 - 3, 2.96, 6, 0.35, 'I N V E S T M E N T',
             size=11, color=SKY, align=PP_ALIGN.CENTER, letter_spacing=10)
    # Location label
    loc = (data.get('title') or '').split(':', 1)[0].strip()
    add_text(slide, MARGIN_X, 4.0, SLIDE_W - 2 * MARGIN_X, 0.4, loc,
             size=12, bold=True, color=SKY, align=PP_ALIGN.CENTER, letter_spacing=4)
    # Title (asset name) — the part after ':'
    full = data.get('title', '')
    name = full.split(':', 1)[1].strip() if ':' in full else full
    add_text(slide, 1.5, 4.45, SLIDE_W - 3, 1.2, name,
             size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Price box
    price = data.get('price', '')
    if price:
        bx = SLIDE_W / 2 - 2.5
        by = 5.7
        add_rect(slide, bx, by, 5.0, 0.95, fill=NAVY, line=SKY, line_width=1)
        add_text(slide, bx, by + 0.08, 5.0, 0.3,
                 (data.get('price_label_total') or 'PRIX TOTAL :'),
                 size=8.5, color=SKY, align=PP_ALIGN.CENTER, letter_spacing=3)
        add_text(slide, bx, by + 0.38, 5.0, 0.55, price,
                 size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Footer
    add_text(slide, MARGIN_X, SLIDE_H - 0.45, SLIDE_W - 2 * MARGIN_X, 0.3,
             'rodschinson.com  |  +32 2 550 36 87  |  assets@rodschinson.com',
             size=8.5, color=GREY_LINE, align=PP_ALIGN.CENTER, letter_spacing=2)


def build_activa(slide, data):
    page_chrome(slide, data, 'a',
                (data.get('section_labels') or {}).get('activa', "Présentation de l'actif"))
    # Two-column body
    left_x, left_w = MARGIN_X, 6.7
    right_x, right_w = left_x + left_w + 0.2, SLIDE_W - (left_x + left_w + 0.2) - MARGIN_X
    top = BODY_TOP + 0.55
    bottom = BODY_BOTTOM

    # Left white card
    add_rect(slide, left_x, top, left_w, bottom - top, fill=WHITE,
             line=GREY_LINE, line_width=0.5)

    full = data.get('title', '')
    title_inside = full  # keep full title here
    add_text(slide, left_x + 0.35, top + 0.25, left_w - 0.7, 0.95, title_inside,
             size=14, bold=True, color=NAVY)

    # Address
    y = top + 1.35
    add_text(slide, left_x + 0.35, y, left_w - 0.7, 0.25,
             data.get('address_label', 'Adresse :'),
             size=9, color=GREY_TX)
    add_text(slide, left_x + 0.35, y + 0.28, left_w - 0.7, 0.3,
             data.get('address', ''), size=11, color=DARK_TX, bold=False)

    # Description bullets
    y = top + 2.05
    add_text(slide, left_x + 0.35, y, left_w - 0.7, 0.25,
             data.get('description_label', 'Description :'),
             size=9, color=GREY_TX)

    # Pull bullets from the description (split on '•')
    desc = data.get('description', '')
    parts = [b.strip() for b in desc.split('•') if b.strip()]
    # drop the lead "BRUXELLES (Etterbeek...)" sentence so we just get bullets
    if parts and ':' in parts[0]:
        parts = parts[1:]
    if not parts:
        parts = [desc]

    bullets_box = slide.shapes.add_textbox(
        Inches(left_x + 0.4), Inches(y + 0.32),
        Inches(left_w - 0.8), Inches(bottom - (y + 0.32) - 0.8))
    tf = bullets_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    for i, b in enumerate(parts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        r = p.add_run()
        r.text = '•  ' + b
        r.font.name = 'Calibri'
        r.font.size = Pt(9.5)
        r.font.color.rgb = DARK_TX
        p.space_after = Pt(3)

    # Price footer
    price = data.get('price', '')
    if price:
        add_rect(slide, left_x + 0.35, bottom - 0.7, left_w - 0.7, 0.012, fill=GREY_LINE)
        add_text(slide, left_x + 0.35, bottom - 0.55, 3.0, 0.4,
                 'Prix total :', size=11, bold=True, color=NAVY)
        add_text(slide, left_x + 1.6, bottom - 0.6, 5.0, 0.5,
                 price, size=18, bold=True, color=SKY)

    # Right photo
    photo = file_path_from_url(data.get('activa_photo') or
                               (data.get('photos') or [''])[0])
    add_picture(slide, photo, right_x, top, right_w, bottom - top)


def build_details_kpi_financial(slide, data):
    page_chrome(slide, data, 'a',
                (data.get('section_labels') or {}).get('details', 'Détails financiers & techniques'))
    add_section_title(slide, MARGIN_X, BODY_TOP + 0.45,
                      SLIDE_W - 2 * MARGIN_X, 0.4,
                      "Détails de l'actif")

    # KPI strip
    metrics = data.get('key_metrics') or []
    if metrics:
        y = BODY_TOP + 1.1
        h = 0.95
        gap = 0.12
        n = len(metrics)
        avail = SLIDE_W - 2 * MARGIN_X
        w = (avail - gap * (n - 1)) / n
        for i, m in enumerate(metrics):
            x = MARGIN_X + i * (w + gap)
            add_rect(slide, x, y, w, h, fill=WHITE, line=GREY_LINE, line_width=0.5)
            # cyan left bar
            add_rect(slide, x, y, 0.045, h, fill=SKY)
            add_text(slide, x + 0.18, y + 0.1, w - 0.25, 0.22,
                     (m.get('label') or '').upper(),
                     size=7.5, bold=True, color=GREY_TX, letter_spacing=4)
            add_text(slide, x + 0.18, y + 0.34, w - 0.25, 0.42,
                     m.get('value') or '',
                     size=16, bold=True, color=NAVY)
            sub = m.get('sub') or ''
            if sub:
                add_text(slide, x + 0.18, y + 0.72, w - 0.25, 0.2,
                         sub, size=7.5, color=GREY_TX)

    # Two columns: financial summary (left), valuation (right)
    cards_top = BODY_TOP + 2.25
    cards_h = BODY_BOTTOM - cards_top - 0.05
    col_w = (SLIDE_W - 2 * MARGIN_X - 0.25) / 2

    draw_data_card(slide, MARGIN_X, cards_top, col_w, cards_h,
                   (data.get('details_labels') or {}).get('financials', 'Résumé financier'),
                   data.get('financial_summary_rows') or [])
    draw_data_card(slide, MARGIN_X + col_w + 0.25, cards_top, col_w, cards_h,
                   (data.get('details_labels') or {}).get('valuation', 'Évaluation'),
                   data.get('valuation_rows') or [])


def draw_data_card(slide, x, y, w, h, head, rows, max_rows=None, head_tag=''):
    """A white card with a navy heading and 2-column rows (label | value)."""
    add_rect(slide, x, y, w, h, fill=WHITE, line=GREY_LINE, line_width=0.5)
    # Heading
    add_text(slide, x + 0.22, y + 0.18, w - 1.0, 0.3, head.upper(),
             size=10, bold=True, color=NAVY, letter_spacing=3)
    if head_tag:
        add_text(slide, x + w - 1.05, y + 0.2, 0.85, 0.25, head_tag,
                 size=8, color=GREY_TX, align=PP_ALIGN.RIGHT)
    # underline accent (sky)
    add_rect(slide, x + 0.22, y + 0.52, 0.6, 0.022, fill=SKY)
    # Rows
    if max_rows is not None:
        rows = rows[:max_rows]
    if not rows:
        return
    row_h = max(0.22, min(0.36, (h - 0.75) / max(len(rows), 1)))
    label_w = (w - 0.45) * 0.6
    val_w   = (w - 0.45) * 0.4
    ry = y + 0.7
    for i, r in enumerate(rows):
        if ry + row_h > y + h - 0.15:
            break
        add_text(slide, x + 0.22, ry, label_w, row_h,
                 r.get('label', ''), size=8.5, color=GREY_TX)
        add_text(slide, x + 0.22 + label_w, ry, val_w, row_h,
                 r.get('value', ''),
                 size=9, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)
        # divider line
        if i < len(rows) - 1 and ry + row_h <= y + h - 0.2:
            add_rect(slide, x + 0.22, ry + row_h - 0.02,
                     w - 0.45, 0.005, fill=GREY_LINE)
        ry += row_h


def build_details_specs(slide, data):
    page_chrome(slide, data, 'a',
                (data.get('section_labels') or {}).get('details', 'Détails financiers & techniques'))
    add_section_title(slide, MARGIN_X, BODY_TOP + 0.45,
                      SLIDE_W - 2 * MARGIN_X, 0.4,
                      "Détails de l'actif", suffix='(SUITE)')
    top = BODY_TOP + 1.1
    h = BODY_BOTTOM - top - 0.05
    draw_data_card(slide, MARGIN_X, top, SLIDE_W - 2 * MARGIN_X, h,
                   (data.get('details_labels') or {}).get('specs', 'Spécifications techniques'),
                   data.get('technical_specs_rows') or [])


def build_unit_table(slide, data):
    page_chrome(slide, data, 'a',
                (data.get('section_labels') or {}).get('details', 'Détails financiers & techniques'))
    add_section_title(slide, MARGIN_X, BODY_TOP + 0.45,
                      SLIDE_W - 2 * MARGIN_X, 0.4,
                      "Détails de l'actif", suffix='(SUITE)')

    top = BODY_TOP + 1.1
    bottom = BODY_BOTTOM
    add_rect(slide, MARGIN_X, top, SLIDE_W - 2 * MARGIN_X, bottom - top,
             fill=WHITE, line=GREY_LINE, line_width=0.5)
    add_text(slide, MARGIN_X + 0.22, top + 0.18, 5.0, 0.3,
             (data.get('details_labels') or {}).get('income', 'Revenus locatifs').upper(),
             size=10, bold=True, color=NAVY, letter_spacing=3)
    add_rect(slide, MARGIN_X + 0.22, top + 0.52, 0.6, 0.022, fill=SKY)

    # Summary row: SURFACE TOTALE, etc.
    summary = []
    for s in data.get('surfaces', []) or []:
        lbl = s.get('floor', s.get('label', ''))
        val = s.get('area', s.get('value', ''))
        if not lbl:
            continue
        if not lbl.lower().startswith(('unité', 'unite', 'studio', 'appart', 'apt', 'lot')):
            summary.append((lbl, val))
    sy = top + 0.7
    sx = MARGIN_X + 0.22
    for lbl, val in summary[:4]:
        add_text(slide, sx, sy, 2.3, 0.22, lbl.upper(),
                 size=7.5, bold=True, color=GREY_TX, letter_spacing=3)
        add_text(slide, sx, sy + 0.24, 2.3, 0.3, val,
                 size=11, bold=True, color=NAVY)
        sx += 2.5
    table_top = sy + 0.7
    add_rect(slide, MARGIN_X + 0.22, table_top - 0.05,
             SLIDE_W - 2 * MARGIN_X - 0.44, 0.005, fill=GREY_LINE)

    # Build unit rows from rental_income_rows pivoted by unit
    units = build_units_from_rows(data.get('rental_income_rows') or [])
    if not units['units']:
        return

    # Header row
    cols = ['Unité', 'Surface'] + units['periods']
    n_periods = len(units['periods'])
    # widths
    unit_w = 2.4
    surf_w = 1.1
    period_w = (SLIDE_W - 2 * MARGIN_X - 0.5 - unit_w - surf_w) / max(n_periods, 1)
    hx = MARGIN_X + 0.22
    hy = table_top + 0.05
    add_text(slide, hx, hy, unit_w, 0.3, 'UNITÉ',
             size=8, bold=True, color=GREY_TX, letter_spacing=3)
    add_text(slide, hx + unit_w, hy, surf_w, 0.3, 'SURFACE',
             size=8, bold=True, color=GREY_TX, letter_spacing=3, align=PP_ALIGN.RIGHT)
    for i, p in enumerate(units['periods']):
        add_text(slide, hx + unit_w + surf_w + i * period_w, hy,
                 period_w, 0.3, p.upper(),
                 size=8, bold=True, color=GREY_TX, letter_spacing=3,
                 align=PP_ALIGN.RIGHT)
    add_rect(slide, hx, hy + 0.36, SLIDE_W - 2 * MARGIN_X - 0.5,
             0.018, fill=NAVY)

    # Data rows
    ry = hy + 0.45
    row_h = (bottom - ry - 0.2) / max(len(units['units']), 1)
    row_h = max(0.25, min(0.35, row_h))
    for i, u in enumerate(units['units']):
        if ry + row_h > bottom - 0.1:
            break
        # zebra stripe
        if i % 2 == 0:
            add_rect(slide, hx - 0.05, ry - 0.03,
                     SLIDE_W - 2 * MARGIN_X - 0.4, row_h,
                     fill=RGBColor(0xFA, 0xFB, 0xFC))
        # unit name + sub
        add_text(slide, hx, ry, unit_w, 0.22, u['unit'],
                 size=9, bold=True, color=DARK_TX)
        if u.get('sub'):
            add_text(slide, hx, ry + 0.18, unit_w, 0.18, u['sub'],
                     size=7.5, color=GREY_TX)
        add_text(slide, hx + unit_w, ry, surf_w, 0.22, u.get('surface', '—'),
                 size=9, color=DARK_TX, align=PP_ALIGN.RIGHT)
        has_step = any(u['steps'].get(p) for p in units['periods'])
        if has_step:
            for j, p in enumerate(units['periods']):
                add_text(slide, hx + unit_w + surf_w + j * period_w, ry,
                         period_w, 0.22, u['steps'].get(p, '—'),
                         size=9, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)
        else:
            # spanned note cell
            add_text(slide, hx + unit_w + surf_w, ry,
                     n_periods * period_w, 0.22,
                     u.get('note', '—'),
                     size=8.5, italic=True, color=GREY_TX,
                     align=PP_ALIGN.RIGHT)
        # row divider
        add_rect(slide, hx, ry + row_h - 0.03,
                 SLIDE_W - 2 * MARGIN_X - 0.4, 0.004, fill=GREY_LINE)
        ry += row_h


def build_units_from_rows(rows):
    import re
    period_kw = re.compile(r'^(mois|month|après|after|année|year|jaar|na|phase|step)\b', re.I)
    units = {}
    order_seq = 0
    periods = []
    for row in rows:
        label = (row.get('label') or '').strip()
        value = (row.get('value') or '').strip()
        if not label or not value:
            continue
        idx = label.rfind(' — ')
        if idx < 0:
            continue
        unit_raw = label[:idx].strip()
        period = label[idx + 3:].strip()
        has_step = bool(period_kw.match(period))
        # surface (… 94 m² …)
        m = re.search(r'\(([^)]*?\d[\d.,]*\s*m²[^)]*)\)', unit_raw, re.I)
        surface = ''
        unit_name = unit_raw
        if m:
            surface = m.group(1).strip()
            unit_name = unit_raw.replace(m.group(0), '').strip().rstrip(',').strip()
        # primary / sub on dash
        primary = unit_name
        sub = ''
        m2 = re.match(r'^(.+?)\s+[–-]\s+(.+)$', unit_name)
        if m2:
            primary, sub = m2.group(1).strip(), m2.group(2).strip()
        key = re.sub(r'[^a-z0-9àâçéèêëîïôûùüÿñæœ]+', ' ',
                     unit_name.lower()).strip()
        u = units.get(key)
        if u is None:
            u = {'unit': primary, 'sub': sub, 'surface': surface,
                 'steps': {}, 'note': '', 'order': order_seq}
            units[key] = u
            order_seq += 1
        else:
            if not u['surface'] and surface: u['surface'] = surface
            if not u['sub'] and sub: u['sub'] = sub
        if has_step:
            u['steps'][period] = value
            if period not in periods:
                periods.append(period)
        elif not u['note']:
            u['note'] = period + (f' — {value}' if value else '')
    ordered = sorted(units.values(), key=lambda x: x['order'])
    return {'units': ordered, 'periods': periods}


def build_details_extras(slide, data):
    page_chrome(slide, data, 'a',
                (data.get('section_labels') or {}).get('details', 'Détails financiers & techniques'))
    add_section_title(slide, MARGIN_X, BODY_TOP + 0.45,
                      SLIDE_W - 2 * MARGIN_X, 0.4,
                      "Détails de l'actif", suffix='(SUITE)')
    top = BODY_TOP + 1.1
    h = BODY_BOTTOM - top - 0.05
    col_w = (SLIDE_W - 2 * MARGIN_X - 0.25) / 2

    # Lease terms left
    draw_data_card(slide, MARGIN_X, top, col_w, h,
                   (data.get('details_labels') or {}).get('leases', 'Conditions du bail'),
                   data.get('lease_terms_rows') or [])

    # Extras bullets right
    bullets_x = MARGIN_X + col_w + 0.25
    add_rect(slide, bullets_x, top, col_w, h, fill=WHITE,
             line=GREY_LINE, line_width=0.5)
    add_text(slide, bullets_x + 0.22, top + 0.18, col_w - 0.45, 0.3,
             (data.get('details_labels') or {}).get('other', 'Autres').upper(),
             size=10, bold=True, color=NAVY, letter_spacing=3)
    add_rect(slide, bullets_x + 0.22, top + 0.52, 0.6, 0.022, fill=SKY)
    items = data.get('extra_bullets') or []
    tb = slide.shapes.add_textbox(
        Inches(bullets_x + 0.22), Inches(top + 0.7),
        Inches(col_w - 0.45), Inches(h - 0.85))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    for i, b in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = '•  ' + b
        r.font.name = 'Calibri'
        r.font.size = Pt(9.5)
        r.font.color.rgb = DARK_TX
        p.space_after = Pt(6)


def build_localisation(slide, data):
    page_chrome(slide, data, 'l',
                (data.get('section_labels') or {}).get('loc', 'Localisation'))
    top = BODY_TOP + 0.55
    bottom = BODY_BOTTOM
    map_w = 7.5
    map_path = file_path_from_url(data.get('map_url'))
    add_picture(slide, map_path, MARGIN_X, top, map_w, bottom - top)

    # Side info panel
    sx = MARGIN_X + map_w + 0.3
    sw = SLIDE_W - sx - MARGIN_X
    add_text(slide, sx, top, sw, 0.4,
             (data.get('section_labels') or {}).get('loc', 'Localisation'),
             size=9, color=SKY, bold=True, letter_spacing=4)
    full = data.get('title', '')
    name = full.split(':', 1)[1].strip() if ':' in full else full
    add_text(slide, sx, top + 0.4, sw, 1.4, name,
             size=14, bold=True, color=NAVY)
    add_rect(slide, sx, top + 1.65, 0.4, 0.018, fill=SKY)

    rows = [
        (data.get('address_label', 'Adresse'), data.get('address', '')),
        ('Référence', data.get('reference', '')),
        (data.get('price_label_total', 'Prix total :').rstrip(' :'),
         data.get('price', '')),
    ]
    ry = top + 1.9
    for k, v in rows:
        if not v: continue
        add_text(slide, sx, ry, sw, 0.22, k.upper(),
                 size=7.5, color=GREY_TX, bold=True, letter_spacing=3)
        add_text(slide, sx, ry + 0.22, sw, 0.3, v,
                 size=10, bold=True, color=DARK_TX)
        ry += 0.65


def build_aerial(slide, data):
    page_chrome(slide, data, 'ae',
                (data.get('section_labels') or {}).get('aerial', 'Vue aérienne'))
    top = BODY_TOP + 0.55
    bottom = BODY_BOTTOM
    map_w = 7.5
    aerial = file_path_from_url(data.get('aerial_view') or data.get('map_url'))
    add_picture(slide, aerial, MARGIN_X, top, map_w, bottom - top)
    # Boundary caption strip
    cap = data.get('boundary_caption', 'Limite indicative de la propriété')
    add_rect(slide, MARGIN_X + 0.15, bottom - 0.55, 3.6, 0.4,
             fill=WHITE, line=SKY, line_width=1.25)
    add_text(slide, MARGIN_X + 0.25, bottom - 0.47, 3.4, 0.3, cap.upper(),
             size=8.5, bold=True, color=NAVY, letter_spacing=3)

    # Side info panel (similar to localisation)
    sx = MARGIN_X + map_w + 0.3
    sw = SLIDE_W - sx - MARGIN_X
    add_text(slide, sx, top, sw, 0.4,
             (data.get('section_labels') or {}).get('aerial', 'Vue aérienne'),
             size=9, color=SKY, bold=True, letter_spacing=4)
    full = data.get('title', '')
    name = full.split(':', 1)[1].strip() if ':' in full else full
    add_text(slide, sx, top + 0.4, sw, 1.4, name,
             size=14, bold=True, color=NAVY)
    add_rect(slide, sx, top + 1.65, 0.4, 0.018, fill=SKY)
    rows = [
        (data.get('address_label', 'Adresse'), data.get('address', '')),
        ('Référence', data.get('reference', '')),
        (data.get('price_label_total', 'Prix total :').rstrip(' :'),
         data.get('price', '')),
    ]
    ry = top + 1.9
    for k, v in rows:
        if not v: continue
        add_text(slide, sx, ry, sw, 0.22, k.upper(),
                 size=7.5, color=GREY_TX, bold=True, letter_spacing=3)
        add_text(slide, sx, ry + 0.22, sw, 0.3, v,
                 size=10, bold=True, color=DARK_TX)
        ry += 0.65


def build_gallery(slide, data, photos):
    page_chrome(slide, data, 'p',
                (data.get('section_labels') or {}).get('photos', 'Reportage photo'))
    top = BODY_TOP + 0.55
    bottom = BODY_BOTTOM
    avail_h = bottom - top
    avail_w = SLIDE_W - 2 * MARGIN_X
    n = len(photos)
    if n <= 3:
        # 1 row, n columns
        cols, rows = n, 1
    elif n <= 4:
        cols, rows = 2, 2
    elif n <= 6:
        cols, rows = 3, 2
    else:
        cols, rows = 3, 3
    gap = 0.12
    cw = (avail_w - gap * (cols - 1)) / cols
    ch = (avail_h - gap * (rows - 1)) / rows
    for i, ph in enumerate(photos):
        r, c = divmod(i, cols)
        x = MARGIN_X + c * (cw + gap)
        y = top + r * (ch + gap)
        add_picture(slide, file_path_from_url(ph), x, y, cw, ch)


def build_sales(slide, data):
    page_chrome(slide, data, 's',
                (data.get('section_labels') or {}).get('sales', 'Conditions de vente'))
    top = BODY_TOP + 0.55
    bottom = BODY_BOTTOM
    left_w = 4.8
    right_x = MARGIN_X + left_w + 0.0
    right_w = SLIDE_W - right_x - MARGIN_X

    # Left card
    add_rect(slide, MARGIN_X, top, left_w, bottom - top, fill=WHITE,
             line=GREY_LINE, line_width=0.5)
    add_text(slide, MARGIN_X + 0.25, top + 0.5, left_w - 0.5, 0.55,
             data.get('agent_name', ''), size=18, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER)
    add_rect(slide, MARGIN_X + 0.25, top + 1.15,
             left_w - 0.5, 0.008, fill=GREY_LINE)
    add_text(slide, MARGIN_X + 0.25, top + 1.3, left_w - 0.5, 0.3,
             data.get('agent_role', ''), size=10, color=DARK_TX,
             align=PP_ALIGN.CENTER)
    add_rect(slide, MARGIN_X + 0.25, top + 1.75,
             left_w - 0.5, 0.008, fill=GREY_LINE)
    add_text(slide, MARGIN_X + 0.4, top + 1.95, left_w - 0.8, 0.3,
             data.get('infos_label', 'Infos :'),
             size=9, color=GREY_TX)
    add_text(slide, MARGIN_X + 0.4, top + 2.3, left_w - 0.8, 0.3,
             data.get('agent_phone', ''),
             size=11, bold=True, color=SKY)
    add_text(slide, MARGIN_X + 0.4, top + 2.65, left_w - 0.8, 0.3,
             data.get('agent_email', ''),
             size=11, bold=True, color=SKY)
    # Price at bottom
    add_rect(slide, MARGIN_X + 0.25, bottom - 0.9,
             left_w - 0.5, 0.008, fill=GREY_LINE)
    add_text(slide, MARGIN_X + 0.4, bottom - 0.75, 2.0, 0.4,
             'Prix total :', size=11, bold=True, color=NAVY)
    add_text(slide, MARGIN_X + 1.65, bottom - 0.8, left_w - 1.8, 0.5,
             data.get('price', ''), size=18, bold=True, color=SKY)

    # Right photo + cta
    photo = file_path_from_url(data.get('sales_photo') or
                               (data.get('photos') or [''])[0])
    add_picture(slide, photo, right_x, top, right_w, (bottom - top) - 0.95)
    # CTA bar
    cta_y = bottom - 0.8
    add_rect(slide, right_x, cta_y, right_w, 0.7,
             fill=WHITE, line=GREY_LINE, line_width=0.5)
    add_text(slide, right_x + 0.15, cta_y + 0.05, right_w - 0.3, 0.22,
             data.get('docs_helper', ''),
             size=8.5, color=GREY_TX, align=PP_ALIGN.CENTER)
    add_rect(slide, right_x + 0.6, cta_y + 0.3, right_w - 1.2, 0.32,
             fill=SKY, line=SKY)
    add_text(slide, right_x + 0.6, cta_y + 0.34, right_w - 1.2, 0.28,
             data.get('docs_label', 'Documents complémentaires'),
             size=9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ─── Main ──────────────────────────────────────────────────────────────────
def build(script_path, output_path):
    with open(script_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]

    def add():
        return prs.slides.add_slide(blank)

    # 1. Cover
    build_cover(add(), data)

    # 2. Activa / Présentation de l'actif
    build_activa(add(), data)

    # 3. Details: KPI + financial summary + valuation
    build_details_kpi_financial(add(), data)

    # 4. Details: technical specs
    build_details_specs(add(), data)

    # 5. Details: unit table (consolidated rental income + surfaces)
    build_unit_table(add(), data)

    # 6. Details: lease terms + extras bullets
    build_details_extras(add(), data)

    # 7. Localisation
    if data.get('map_url'):
        build_localisation(add(), data)

    # 8. Vue aérienne
    if data.get('aerial_view') or data.get('map_url'):
        build_aerial(add(), data)

    # 9+. Photos gallery — chunk to ~6 per slide
    photos = [p for p in (data.get('photos') or []) if p]
    if photos:
        i = 0
        # First slide: hero trio if we have ≥ 3
        if len(photos) >= 3:
            build_gallery(add(), data, photos[:3])
            i = 3
        # Then 6-up slides
        while i < len(photos):
            n = min(6, len(photos) - i)
            build_gallery(add(), data, photos[i:i + n])
            i += n

    # Last: Conditions de vente / Bea Neetens contact
    build_sales(add(), data)

    prs.save(output_path)
    print(f'PPTX → {output_path}')
    print(f'  slides: {len(prs.slides)}, size: {os.path.getsize(output_path):,} bytes')


if __name__ == '__main__':
    ap = argparse.ArgumentParser()
    ap.add_argument('--script', required=True)
    ap.add_argument('--output', required=True)
    args = ap.parse_args()
    build(args.script, args.output)
